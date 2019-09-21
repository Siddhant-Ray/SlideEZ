from pptx import Presentation
from pptx.util import Inches, Pt 

presentation = "testppt3.pptx"
prs = Presentation(pptx=None)
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Test"
subtitle.text = "test"
prs.save(presentation)
text_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(text_slide_layout)
title = slide.shapes.title
title.text= slide_title
content = slide.shapes.placeholders[1]
tf = content.text_frame
p=tf.add_paragraph()
p.text=slide_subtopic
p.level=1
prs.save(presentation)
image_slide_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(image_slide_layout)
title = slide.shapes.title
title.text=slide_title
content = slide.shapes.placeholders[1]
tf = content.text_frame
p=tf.add_paragraph()
p.text=slide_subtopic
p.level=1
left = Inches(6)
top = Inches(3)
height = Inches(2)
pic = slide.shapes.add_picture(img_path, left, top, height=height)
prs.save(presentation)