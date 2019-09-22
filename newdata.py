
def get_result():
    import pandas
    import re
    import nltk
    nltk.download('stopwords')
    from nltk.corpus import stopwords
    from nltk.stem.porter import PorterStemmer
    from nltk.tokenize import RegexpTokenizer
    from urllib.request import urlretrieve
    from urllib.parse import quote 
    

    nltk.download('wordnet')
    from nltk.stem.wordnet import WordNetLemmatizer

    f=open("keyword.txt",'r')

    dataset=f.read()
    f.close()
    print(dataset)

    corpus=[]
    ##Creating a list of stop words and adding custom stopwords
    stop_words = set(stopwords.words("english"))##Creating a list of custom stopwords
    new_words = ["using", "show", "result", "large", "also", "iv", "one", "two", "new", "previously", "shown"]
    stop_words = stop_words.union(new_words)

    test=dataset
    test=test.lower()
    test=test.split()
    ps=PorterStemmer()
    lem=WordNetLemmatizer()
    test=[lem.lemmatize(word) for word in test if not word in stop_words]
    test=' '.join(test)
    print(test)

    corpus.append(test)
    stop_words = set(stopwords.words("english"))##Creating a list of custom stopwords
    new_words = ["using", "show", "result", "large", "also", "iv", "one", "two", "new", "previously", "shown"]
    stop_words = stop_words.union(new_words)


    from sklearn.feature_extraction.text import CountVectorizer
    import re


    cv=CountVectorizer(stop_words=stop_words,ngram_range=(1,3))
    X=cv.fit_transform(corpus)


    #Most frequently occuring words
    def get_top_n_words(corpus, n=None):
        vec = CountVectorizer().fit(corpus)
        bag_of_words = vec.transform(corpus)
        sum_words = bag_of_words.sum(axis=0)
        words_freq = [(word, sum_words[0, idx]) for word, idx in
                       vec.vocabulary_.items()]
        words_freq =sorted(words_freq, key = lambda x: x[1],
                           reverse=True)
        return words_freq[:n]#Convert most freq words to dataframe for plotting bar plot
    top_words = get_top_n_words(corpus, n=4)
    top_df = pandas.DataFrame(top_words)
    
    top_df.columns=["Word", "Freq"]#Barplot of most freq words

    unigram=list(top_df["Word"][:4])
    print(unigram)
    print(top_df["Word"][:4])
    print(top_df)

    #Most frequently occuring Bi-grams
    def get_top_n2_words(corpus, n=None):
        vec1 = CountVectorizer(ngram_range=(2,2),
                max_features=2000).fit(corpus)
        bag_of_words = vec1.transform(corpus)
        sum_words = bag_of_words.sum(axis=0)
        words_freq = [(word, sum_words[0, idx]) for word, idx in
                      vec1.vocabulary_.items()]
        words_freq =sorted(words_freq, key = lambda x: x[1],
                    reverse=True)
        return words_freq[:n]


    top2_words=get_top_n2_words(corpus, n=4)
    top2_df = pandas.DataFrame(top2_words)
    top2_df.columns=["Bi-gram", "Freq"]
    print(top2_df)#Barplot of most freq Bi-grams


    #Most frequently occuring Tri-grams
    def get_top_n3_words(corpus, n=None):
        vec1 = CountVectorizer(ngram_range=(3,3),
               max_features=2000).fit(corpus)
        bag_of_words = vec1.transform(corpus)
        sum_words = bag_of_words.sum(axis=0)
        words_freq = [(word, sum_words[0, idx]) for word, idx in
                      vec1.vocabulary_.items()]
        words_freq =sorted(words_freq, key = lambda x: x[1],
                    reverse=True)
        return words_freq[:n]


    top3_words = get_top_n3_words(corpus, n=4)
    top3_df = pandas.DataFrame(top3_words)
    top3_df.columns=["Tri-gram", "Freq"]
    print(top3_df)#Barplot of most freq Tri-grams

    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from PIL import Image
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
    
    prs=Presentation("testppt3.pptx")
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    #shape=slide.shapes[0]
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 40
    gradient_stops = fill.gradient_stops
    gradient_stop = gradient_stops[0]
    color = gradient_stop.color
    color.theme_color = MSO_THEME_COLOR.LIGHT_1

    #text_frame = shape.text_frame
    #text_frame.clear()
    content = slide.shapes.placeholders[1]
    text_frame = content.text_frame
    title=slide.shapes.title
    title.text="Keywords"
    
    

    for i in range(len(unigram)-1):
        p = text_frame.paragraphs[0]
        run=p.add_run()
        run.text = unigram[i]+"\n"
        run.level=1
        run.hyperlink.address = 'https://www.duckduckgo.com/?q='+unigram[i]
        
    p = text_frame.paragraphs[0]
    run=p.add_run()
    run.text = unigram[-1]
    run.level=1
    run.hyperlink.address = 'https://www.duckduckgo.com/?q='+unigram[-1]
    prs.save("testppt3.pptx")
#get_result()

