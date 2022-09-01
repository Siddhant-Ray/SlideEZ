def fn(test):
    
    from deepsegment import DeepSegment
    segmenter=DeepSegment('en')
    import textrazor
    textrazor.api_key = "043e170ef41a6d297a508581225bd493943f3a9f831345fb71f86d64"

    client = textrazor.TextRazor(extractors=["words", "relations"])
    #client.set_do_cleanup_HTML(True)

    response = client.analyze(test)
    l=[]

    for property in response.properties():
        for word in property.predicate_words:
            l.append(word.lemma)
            if word.lemma == "sound":
                for property_word in property.property_words:
                    for phrase in property_word.noun_phrases:
                        print (phrase)
                break
    l=[]
    flag=False
    for sentence in response.sentences():
        print(sentence.words)
        for word in sentence.words:
            if word.lemma=="image" or word.lemma=="picture" or word.lemma=="photo" or word.lemma=="show" or word.lemma=="see" or word.lemma=="display":
                k=word.lemma
                flag=True 
            l.append(word.lemma)
    astring=""
    for i in l:
        astring+=i+" "

    f=open("keyword.txt",'a')
    f.write(astring+"\n")
    f.close()
    alist=segmenter.segment(astring)
    print(alist)

    if(flag):
        s=l.index(k)
        m=l[s:]
    
        t=""
        st=""
        for i in m:
            t+=i+" "
    else:
        t="No image found"
        st=""
        for j in l:
            st+=j+" "

    text1=st
    text2=t

    print(t)
    response1=client.analyze(t)

    for noun in response1.noun_phrases():
        print(noun.words)
        for word in noun.words:
            print(word.lemma)


    from requests import exceptions
    import argparse
    import requests
    import cv2
    import os
    import time

    starttime=time.time();
    

    # set your Microsoft Cognitive Services API key along with (1) the
    # maximum number of results for a given search and (2) the group size
    # for results (maximum of 50 per request)
    API_KEY = "1539329da5c84f4aa90a8789509db5a9"
    MAX_RESULTS = 1
    GROUP_SIZE = 1
    
    # set the endpoint API URL
    URL = "https://api.cognitive.microsoft.com/bing/v7.0/images/search"

    # when attempting to download images from the web both the Python
    # programming language and the requests library have a number of
    # exceptions that can be thrown so let's build a list of them now
    # so we can filter on them
    EXCEPTIONS = set([IOError, FileNotFoundError,
        exceptions.RequestException, exceptions.HTTPError,
        exceptions.ConnectionError, exceptions.Timeout])


    # store the search term in a convenience variable then set the
    # headers and search parameters
    term = t
    headers = {"Ocp-Apim-Subscription-Key" : API_KEY}
    params = {"q": term, "offset": 0, "count": GROUP_SIZE}
    
    # make the search
    print("[INFO] searching Bing API for '{}'".format(term))
    search = requests.get(URL, headers=headers, params=params)
    search.raise_for_status()
    
    # grab the results from the search, including the total number of
    # estimated results returned by the Bing API
    results = search.json()
    estNumResults = min(results["totalEstimatedMatches"], MAX_RESULTS)
    print("[INFO] {} total results for '{}'".format(estNumResults,
        term))
    
    # initialize the total number of images downloaded thus far
    total = 0


    for offset in range(0, estNumResults, GROUP_SIZE):
        # update the search parameters using the current offset, then
        # make the request to fetch the results
        print("[INFO] making request for group {}-{} of {}...".format(
            offset, offset + GROUP_SIZE, estNumResults))
        params["offset"] = offset
        search = requests.get(URL, headers=headers, params=params)
        search.raise_for_status()
        results = search.json()
        print("[INFO] saving images for group {}-{} of {}...".format(
            offset, offset + GROUP_SIZE, estNumResults))
            # loop over the results
        for v in results["value"]:
            # try to download the image
            try:
                # make a request to download the image
                print("[INFO] fetching: {}".format(v["contentUrl"]))
                r = requests.get(v["contentUrl"], timeout=30)
    
                # build the path to the output image
                ext = v["contentUrl"][v["contentUrl"].rfind("."):]
                p = os.path.sep.join([r'C:\Users\HP\Desktop\Projects\VIT Hack\SlideEZ', "{}{}".format(str(total).zfill(8), ext)])

                print("The answer is")
                print(p)
    
                # write the image to disk
                f = open(p, "wb")
                f.write(r.content)
                f.close()
    
            # catch any errors that would not unable us to download the
            # image
            except Exception as e:
                # check to see if our exception is in our list of
                # exceptions to check for
                if type(e) in EXCEPTIONS:
                    print("[INFO] skipping: {}".format(v["contentUrl"]))
                    continue
            # try to load the image from disk
            image = cv2.imread(p)

            # if the image is `None` then we could not properly load the
            # image from disk (so it should be ignored)
            if image is None:
                print("[INFO] deleting: {}".format(p))
                os.remove(p)
                continue

            # update the counter
            total += 1
    endtime=time.time()-starttime
    print("Total time taken to search for the query is")
    print(endtime)

    from pptx import Presentation
    from pptx.util import Inches, Pt 
    from pptx.enum.text import PP_ALIGN
    from PIL import Image
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR

    presentation = "testppt3.pptx"
    prs = Presentation(presentation)
    if len(prs.slides)==0:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        background=slide.background
        fill=background.fill
        fill.gradient()
        fill.gradient_angle=40
        gradient_stops=fill.gradient_stops
        gradient_stop=gradient_stops[0]
        color=gradient_stop.color
        color.theme_color=MSO_THEME_COLOR.LIGHT_1
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Test"
        subtitle.text = "test"
        prs.save(presentation)
    if not flag:

        text_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(text_slide_layout)
        background=slide.background
        fill=background.fill
        fill.gradient()
        fill.gradient_angle=40
        gradient_stops=fill.gradient_stops
        gradient_stop=gradient_stops[0]
        color=gradient_stop.color
        color.theme_color=MSO_THEME_COLOR.LIGHT_1
        title = slide.shapes.title
        blist=[]
        for i in range(0,len(alist)):
            blist+=alist[i].split(" ")

        mx=0
        slide_t=""
        for j in blist:
            if(len(j)>=mx):
                mx=len(j)
                slide_t=j.title()
            
        title.text= slide_t
        content = slide.shapes.placeholders[1]
        tf = content.text_frame
        for i in alist:
            para=tf.add_paragraph()
            para.text=i
            para.level=1
        prs.save(presentation)
    else:

        image_slide_layout = prs.slide_layouts[8]
        slide = prs.slides.add_slide(image_slide_layout)
        background=slide.background
        fill=background.fill
        fill.gradient()
        fill.gradient_angle=40
        gradient_stops=fill.gradient_stops
        gradient_stop=gradient_stops[0]
        color=gradient_stop.color
        color.theme_color=MSO_THEME_COLOR.LIGHT_1
        #title = slide.shapes.title
        #title.text="Sub2"
        content = slide.shapes.placeholders[1]
        im=Image.open(p)
        width,height= im.size
        content.height= height
        content.width= width
        content.insert_picture(p)
        content = slide.shapes.placeholders[0]
        tf = content.text_frame
        for i in alist:
            
            para=tf.add_paragraph()
            para.text=i
            para.level=1
            para.alignment=PP_ALIGN.CENTER
        #left = Inches(6)
        #top = Inches(3)
        #height = Inches(2)
        #pic = slide.shapes.add_picture(p, left, top, height=height)
        prs.save(presentation)


